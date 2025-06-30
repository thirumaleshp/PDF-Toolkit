import os
import io
from typing import List, Union
from PyPDF2 import PdfWriter, PdfReader
from reportlab.lib.pagesizes import letter, A4
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image as RLImage
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from PIL import Image
import docx
from docx.shared import Inches
import openpyxl
from openpyxl.drawing.image import Image as OpenpyxlImage
from pptx import Presentation
from pptx.util import Inches as PptxInches
import markdown
import img2pdf
import fitz  # PyMuPDF
import tempfile

class PDFToolkit:
    def __init__(self):
        self.supported_formats = {
            'images': ['.jpg', '.jpeg', '.png', '.bmp', '.tiff', '.gif'],
            'documents': ['.docx', '.xlsx', '.pptx', '.txt', '.md'],
            'pdf': ['.pdf']
        }

    def split_pdf(self, pdf_file: str, page_ranges: Union[str, List[int]], output_path: str, combine_pages: bool = True) -> bool:
        """Split a PDF file into specified pages.
        
        Args:
            pdf_file: Path to the input PDF file
            page_ranges: String like "2-8,10,12" or list of page numbers
            output_path: Path for the output PDF file (if combine_pages=True) or directory (if combine_pages=False)
            combine_pages: If True, create one PDF with selected pages. If False, create separate PDFs for each page
        """
        try:
            # Parse page ranges
            if isinstance(page_ranges, str):
                page_list = []
                ranges = page_ranges.replace(' ', '').split(',')
                for item in ranges:
                    if '-' in item:
                        start, end = item.split('-')
                        page_list.extend(list(range(int(start), int(end) + 1)))
                    else:
                        page_list.append(int(item))
            else:
                page_list = page_ranges

            reader = PdfReader(pdf_file)
            total_pages = len(reader.pages)
            
            # Validate page numbers
            valid_pages = []
            for page_num in page_list:
                if 1 <= page_num <= total_pages:
                    valid_pages.append(page_num)
                else:
                    print(f"Warning: Page {page_num} is out of range (1-{total_pages}), skipping...")
            
            if not valid_pages:
                print("Error: No valid pages specified")
                return False

            if combine_pages:
                # Create one PDF with all selected pages
                writer = PdfWriter()
                for page_num in valid_pages:
                    writer.add_page(reader.pages[page_num - 1])
                
                with open(output_path, 'wb') as output_file:
                    writer.write(output_file)
            else:
                # Create separate PDFs for each page
                os.makedirs(output_path, exist_ok=True)
                for page_num in valid_pages:
                    writer = PdfWriter()
                    writer.add_page(reader.pages[page_num - 1])
                    
                    page_output_path = os.path.join(output_path, f"page_{page_num}.pdf")
                    with open(page_output_path, 'wb') as output_file:
                        writer.write(output_file)

            return True
        except Exception as e:
            print(f"Error splitting PDF: {str(e)}")
            return False
    
    def merge_pdfs(self, pdf_files: List[str], output_path: str) -> bool:
        """Merge multiple PDF files into one."""
        try:
            writer = PdfWriter()
            
            for pdf_file in pdf_files:
                if not os.path.exists(pdf_file):
                    print(f"Warning: File {pdf_file} not found, skipping...")
                    continue
                
                reader = PdfReader(pdf_file)
                for page in reader.pages:
                    writer.add_page(page)
            
            with open(output_path, 'wb') as output_file:
                writer.write(output_file)
            
            return True
        except Exception as e:
            print(f"Error merging PDFs: {str(e)}")
            return False
    
    def convert_to_pdf(self, input_file: str, output_path: str) -> bool:
        """Convert various file formats to PDF."""
        try:
            file_ext = os.path.splitext(input_file)[1].lower()
            
            if file_ext in self.supported_formats['images']:
                return self._convert_image_to_pdf(input_file, output_path)
            elif file_ext == '.docx':
                return self._convert_docx_to_pdf(input_file, output_path)
            elif file_ext == '.xlsx':
                return self._convert_xlsx_to_pdf(input_file, output_path)
            elif file_ext == '.pptx':
                return self._convert_pptx_to_pdf(input_file, output_path)
            elif file_ext == '.txt':
                return self._convert_txt_to_pdf(input_file, output_path)
            elif file_ext == '.md':
                return self._convert_md_to_pdf(input_file, output_path)
            else:
                print(f"Unsupported file format: {file_ext}")
                return False
        
        except Exception as e:
            print(f"Error converting file: {str(e)}")
            return False
    
    def _convert_image_to_pdf(self, image_path: str, output_path: str) -> bool:
        """Convert image to PDF using img2pdf for better quality."""
        try:
            with open(image_path, "rb") as image_file:
                pdf_bytes = img2pdf.convert(image_file.read())
            
            with open(output_path, "wb") as pdf_file:
                pdf_file.write(pdf_bytes)
            
            return True
        except Exception as e:
            # Fallback to PIL/reportlab method
            try:
                img = Image.open(image_path)
                if img.mode != 'RGB':
                    img = img.convert('RGB')
                
                doc = SimpleDocTemplate(output_path, pagesize=A4)
                story = []
                
                # Calculate image size to fit page
                img_width, img_height = img.size
                page_width, page_height = A4
                aspect = img_height / float(img_width)
                
                if img_width > page_width - 2*inch:
                    img_width = page_width - 2*inch
                    img_height = img_width * aspect
                
                if img_height > page_height - 2*inch:
                    img_height = page_height - 2*inch
                    img_width = img_height / aspect
                
                story.append(RLImage(image_path, width=img_width, height=img_height))
                doc.build(story)
                return True
            except Exception as e2:
                print(f"Error converting image: {str(e2)}")
                return False
    
    def _convert_docx_to_pdf(self, docx_path: str, output_path: str) -> bool:
        """Convert DOCX to PDF."""
        try:
            doc = docx.Document(docx_path)
            pdf_doc = SimpleDocTemplate(output_path, pagesize=A4)
            styles = getSampleStyleSheet()
            story = []
            
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    story.append(Paragraph(paragraph.text, styles['Normal']))
                    story.append(Spacer(1, 12))
            
            pdf_doc.build(story)
            return True
        except Exception as e:
            print(f"Error converting DOCX: {str(e)}")
            return False
    
    def _convert_xlsx_to_pdf(self, xlsx_path: str, output_path: str) -> bool:
        """Convert XLSX to PDF."""
        try:
            workbook = openpyxl.load_workbook(xlsx_path)
            pdf_doc = SimpleDocTemplate(output_path, pagesize=A4)
            styles = getSampleStyleSheet()
            story = []
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                story.append(Paragraph(f"Sheet: {sheet_name}", styles['Heading1']))
                story.append(Spacer(1, 12))
                
                for row in sheet.iter_rows(values_only=True):
                    row_text = " | ".join([str(cell) if cell is not None else "" for cell in row])
                    if row_text.strip():
                        story.append(Paragraph(row_text, styles['Normal']))
                        story.append(Spacer(1, 6))
                
                story.append(Spacer(1, 24))
            
            pdf_doc.build(story)
            return True
        except Exception as e:
            print(f"Error converting XLSX: {str(e)}")
            return False
    
    def _convert_pptx_to_pdf(self, pptx_path: str, output_path: str) -> bool:
        """Convert PPTX to PDF."""
        try:
            prs = Presentation(pptx_path)
            pdf_doc = SimpleDocTemplate(output_path, pagesize=A4)
            styles = getSampleStyleSheet()
            story = []
            
            for i, slide in enumerate(prs.slides):
                story.append(Paragraph(f"Slide {i+1}", styles['Heading1']))
                story.append(Spacer(1, 12))
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        story.append(Paragraph(shape.text, styles['Normal']))
                        story.append(Spacer(1, 6))
                
                story.append(Spacer(1, 24))
            
            pdf_doc.build(story)
            return True
        except Exception as e:
            print(f"Error converting PPTX: {str(e)}")
            return False
    
    def _convert_txt_to_pdf(self, txt_path: str, output_path: str) -> bool:
        """Convert TXT to PDF."""
        try:
            with open(txt_path, 'r', encoding='utf-8') as file:
                content = file.read()
            
            pdf_doc = SimpleDocTemplate(output_path, pagesize=A4)
            styles = getSampleStyleSheet()
            story = []
            
            for line in content.split('\n'):
                if line.strip():
                    story.append(Paragraph(line, styles['Normal']))
                story.append(Spacer(1, 12))
            
            pdf_doc.build(story)
            return True
        except Exception as e:
            print(f"Error converting TXT: {str(e)}")
            return False
    
    def _convert_md_to_pdf(self, md_path: str, output_path: str) -> bool:
        """Convert Markdown to PDF."""
        try:
            with open(md_path, 'r', encoding='utf-8') as file:
                md_content = file.read()
            
            html_content = markdown.markdown(md_content)
            
            pdf_doc = SimpleDocTemplate(output_path, pagesize=A4)
            styles = getSampleStyleSheet()
            story = []
            
            # Simple HTML to PDF conversion (basic)
            story.append(Paragraph(html_content, styles['Normal']))
            
            pdf_doc.build(story)
            return True
        except Exception as e:
            print(f"Error converting Markdown: {str(e)}")
            return False
    
    def get_supported_formats(self) -> dict:
        """Return supported file formats."""
        return self.supported_formats
    
    def is_supported_format(self, file_path: str) -> bool:
        """Check if file format is supported."""
        file_ext = os.path.splitext(file_path)[1].lower()
        for format_type, extensions in self.supported_formats.items():
            if file_ext in extensions:
                return True
        return False
    
    def convert_pdf_to_word(self, pdf_path: str, output_path: str) -> bool:
        """Convert PDF to Word document."""
        try:
            # Open PDF
            pdf_document = fitz.open(pdf_path)
            doc = docx.Document()
            
            for page_num in range(pdf_document.page_count):
                page = pdf_document[page_num]
                text = page.get_text()
                
                if text.strip():
                    # Add page break for new pages (except first)
                    if page_num > 0:
                        doc.add_page_break()
                    
                    # Split text into paragraphs and add to document
                    paragraphs = text.split('\n\n')
                    for para in paragraphs:
                        if para.strip():
                            doc.add_paragraph(para.strip())
            
            doc.save(output_path)
            pdf_document.close()
            return True
            
        except Exception as e:
            print(f"Error converting PDF to Word: {str(e)}")
            return False
    
    def convert_pdf_to_excel(self, pdf_path: str, output_path: str) -> bool:
        """Convert PDF to Excel spreadsheet."""
        try:
            # Open PDF
            pdf_document = fitz.open(pdf_path)
            workbook = openpyxl.Workbook()
            
            # Remove default sheet
            workbook.remove(workbook.active)
            
            for page_num in range(pdf_document.page_count):
                page = pdf_document[page_num]
                text = page.get_text()
                
                # Create a new sheet for each page
                sheet = workbook.create_sheet(title=f"Page_{page_num + 1}")
                
                if text.strip():
                    lines = text.split('\n')
                    for row_idx, line in enumerate(lines, 1):
                        if line.strip():
                            # Try to split by common delimiters
                            if '\t' in line:
                                cells = line.split('\t')
                            elif '|' in line:
                                cells = line.split('|')
                            elif '  ' in line:  # Multiple spaces
                                cells = [cell.strip() for cell in line.split('  ') if cell.strip()]
                            else:
                                cells = [line.strip()]
                            
                            for col_idx, cell in enumerate(cells, 1):
                                sheet.cell(row=row_idx, column=col_idx, value=cell.strip())
            
            workbook.save(output_path)
            pdf_document.close()
            return True
            
        except Exception as e:
            print(f"Error converting PDF to Excel: {str(e)}")
            return False
    
    def convert_pdf_to_powerpoint(self, pdf_path: str, output_path: str) -> bool:
        """Convert PDF to PowerPoint presentation."""
        try:
            # Open PDF
            pdf_document = fitz.open(pdf_path)
            prs = Presentation()
            
            for page_num in range(pdf_document.page_count):
                page = pdf_document[page_num]
                text = page.get_text()
                
                # Create slide
                slide_layout = prs.slide_layouts[1]  # Title and Content layout
                slide = prs.slides.add_slide(slide_layout)
                
                # Set title
                title = slide.shapes.title
                title.text = f"Page {page_num + 1}"
                
                # Add content
                if text.strip():
                    content = slide.placeholders[1]
                    text_frame = content.text_frame
                    text_frame.text = text.strip()
            
            prs.save(output_path)
            pdf_document.close()
            return True
            
        except Exception as e:
            print(f"Error converting PDF to PowerPoint: {str(e)}")
            return False
    
    def convert_pdf_to_images(self, pdf_path: str, output_dir: str, image_format: str = 'png') -> bool:
        """Convert PDF pages to images."""
        try:
            # Open PDF
            pdf_document = fitz.open(pdf_path)
            os.makedirs(output_dir, exist_ok=True)
            
            for page_num in range(pdf_document.page_count):
                page = pdf_document[page_num]
                
                # Render page to image
                mat = fitz.Matrix(2.0, 2.0)  # Zoom factor for better quality
                pix = page.get_pixmap(matrix=mat)
                
                # Save image
                output_path = os.path.join(output_dir, f"page_{page_num + 1}.{image_format}")
                if image_format.lower() == 'png':
                    pix.save(output_path)
                else:
                    # Convert to PIL Image for other formats
                    img_data = pix.tobytes("ppm")
                    img = Image.open(io.BytesIO(img_data))
                    img.save(output_path, image_format.upper())
            
            pdf_document.close()
            return True
            
        except Exception as e:
            print(f"Error converting PDF to images: {str(e)}")
            return False
    
    def get_pdf_info(self, pdf_path: str) -> dict:
        """Get PDF information including page count."""
        try:
            pdf_document = fitz.open(pdf_path)
            info = {
                'page_count': pdf_document.page_count,
                'metadata': pdf_document.metadata
            }
            pdf_document.close()
            return info
        except Exception as e:
            print(f"Error getting PDF info: {str(e)}")
            return {'page_count': 0, 'metadata': {}}
